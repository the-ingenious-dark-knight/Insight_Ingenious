"""Insight Ingenious – Unstructured Extractor *format* smoke tests
=================================================================

This module provides minimal, end-to-end smoke tests for the UnstructuredExtractor
backend in Insight Ingenious. It verifies that the extractor can handle
common non-PDF rich text formats (DOCX and PPTX) in both byte-stream and
file-path modes. Each test ensures that at least one text element with non-empty
`text` is extracted, confirming basic integration without asserting deep
semantic correctness.

Test scenarios:
  - DOCX bytes: feed raw Microsoft Word content into the extractor
  - DOCX file: extract directly from a `.docx` file on disk
  - PPTX file: generate a minimal PowerPoint presentation via python-pptx
  - PPTX bytes: read that presentation back into a byte array and extract

Skip conditions:
  - All tests are skipped if the `unstructured` library is not installed
  - PPTX-specific tests are also skipped if the `python-pptx` library is not installed

Execution:
  Run this module in isolation from the repository root:

    uv run pytest -q ingenious/document_processing/tests/test_unstructured_formats.py

All assertions should pass without errors, confirming that the extractor
produces at least one non-empty text element for each scenario.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from ingenious.document_processing.extractor import _load

# ──────────────────────────────────────────────────────────────────────────────
# Tests for DOCX input (bytes and file path)
# ──────────────────────────────────────────────────────────────────────────────


def test_unstructured_docx_bytes(
    unstructured_available: bool,
    docx_bytes: bytes,
) -> None:
    """
    Verify the extractor handles raw DOCX bytes correctly.

    Skips if the `unstructured` library is not installed.

    Parameters
    ----------
    unstructured_available : bool
        True if the `unstructured` library is installed; else the test is skipped.
    docx_bytes : bytes
        Raw byte content of a Microsoft Word document (from conftest).

    Assertions
    ----------
    * At least one element is returned by `ux.extract`.
    * At least one element contains non-empty `text`.
    """
    if not unstructured_available:
        pytest.skip("unstructured not installed")

    ux = _load("unstructured")
    elements = list(ux.extract(docx_bytes))

    assert elements, "No elements extracted from DOCX bytes"
    assert any(el["text"] for el in elements), "DOCX extraction produced no text"


def test_unstructured_docx_path(
    unstructured_available: bool,
    docx_path: Path,
) -> None:
    """
    Verify the extractor handles a DOCX file path correctly.

    Skips if the `unstructured` library is not installed.

    Parameters
    ----------
    unstructured_available : bool
        True if the `unstructured` library is installed; else the test is skipped.
    docx_path : pathlib.Path
        File path to a `.docx` document (from conftest).

    Assertions
    ----------
    * At least one element is returned by `ux.extract`.
    * At least one element contains non-empty `text`.
    """
    if not unstructured_available:
        pytest.skip("unstructured not installed")

    ux = _load("unstructured")
    elements = list(ux.extract(docx_path))

    assert elements, "No elements extracted from DOCX file"
    assert any(el["text"] for el in elements), "DOCX file extraction produced no text"


# ──────────────────────────────────────────────────────────────────────────────
# Fixtures and tests for PPTX input (file path and bytes)
# ──────────────────────────────────────────────────────────────────────────────


@pytest.fixture(scope="session")
def pptx_available() -> bool:
    """
    Return True if the optional library `python-pptx` can be imported.
    """
    try:
        import pptx  # noqa: F401

        return True
    except ImportError:
        return False


@pytest.fixture(scope="function")
def pptx_path(tmp_path: Path, pptx_available: bool) -> Path:
    """
    Create a minimal PowerPoint file and return its path.

    Skips the test if `python-pptx` is not installed.

    Parameters
    ----------
    tmp_path : pathlib.Path
        pytest-provided temporary directory unique to this test.
    pptx_available : bool
        True if `python-pptx` is installed; else the test is skipped.

    Returns
    -------
    pathlib.Path
        Path to the generated `sample.pptx` file.
    """
    if not pptx_available:
        pytest.skip("python-pptx not installed")

    from pptx import Presentation  # type: ignore

    output = tmp_path / "sample.pptx"
    pres: Any = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[0])
    slide.shapes.title.text = "Hello PPTX"
    pres.save(output)
    return output


@pytest.fixture(scope="function")
def pptx_bytes(pptx_path: Path) -> bytes:
    """
    Read and return the raw bytes of the PowerPoint file created by `pptx_path`.

    Parameters
    ----------
    pptx_path : pathlib.Path
        Path to a valid `.pptx` presentation.

    Returns
    -------
    bytes
        Binary content of the PowerPoint file.
    """
    return pptx_path.read_bytes()


def test_unstructured_pptx_bytes(
    unstructured_available: bool,
    pptx_available: bool,
    pptx_bytes: bytes,
) -> None:
    """
    Verify the extractor handles raw PPTX bytes correctly.

    Skips if either `unstructured` or `python-pptx` is not installed.

    Parameters
    ----------
    unstructured_available : bool
        True if the `unstructured` library is installed; else the test is skipped.
    pptx_available : bool
        True if the `python-pptx` library is installed; else the test is skipped.
    pptx_bytes : bytes
        Raw byte content of a PowerPoint presentation.

    Assertions
    ----------
    * At least one element is returned by `ux.extract`.
    * At least one element contains non-empty `text`.
    """
    if not (unstructured_available and pptx_available):
        pytest.skip("unstructured or python-pptx missing")

    ux = _load("unstructured")
    elements = list(ux.extract(pptx_bytes))

    assert elements, "No elements extracted from PPTX bytes"
    assert any(el["text"] for el in elements), "PPTX extraction produced no text"


def test_unstructured_pptx_path(
    unstructured_available: bool,
    pptx_available: bool,
    pptx_path: Path,
) -> None:
    """
    Verify the extractor handles a PPTX file path correctly.

    Skips if either `unstructured` or `python-pptx` is not installed.

    Parameters
    ----------
    unstructured_available : bool
        True if the `unstructured` library is installed; else the test is skipped.
    pptx_available : bool
        True if the `python-pptx` library is installed; else the test is skipped.
    pptx_path : pathlib.Path
        Path to a valid `.pptx` presentation.

    Assertions
    ----------
    * At least one element is returned by `ux.extract`.
    * At least one element contains non-empty `text`.
    """
    if not (unstructured_available and pptx_available):
        pytest.skip("unstructured or python-pptx missing")

    ux = _load("unstructured")
    elements = list(ux.extract(pptx_path))

    assert elements, "No elements extracted from PPTX file"
    assert any(el["text"] for el in elements), "PPTX file extraction produced no text"
