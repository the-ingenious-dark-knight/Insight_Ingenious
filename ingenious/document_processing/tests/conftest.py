"""insight-ingenious – shared *pytest* configuration
===================================================

Centralised fixtures and helpers for the document-processing test-suite.

Key design points
-----------------
* **No committed binaries** – sample PDFs download at run-time and cache to
  ``tmp_path``; DOCX and PPTX samples are generated on-the-fly.
* **Graceful network degradation** – network failures *xfail* when a cached
  copy exists, otherwise affected tests are skipped.
* **Optional-dependency discovery** – fixtures detect whether *PyMuPDF*,
  *python-docx*, *pdfminer.six*, *python-pptx*, or *unstructured* are present
  and adjust the matrix automatically.

All public helpers carry full NumPy-style docstrings and are enumerated in
:pydata:`__all__` for static-analysis tools.
"""

from __future__ import annotations

import importlib
import importlib.util
from collections.abc import Callable
from pathlib import Path
from typing import Any, Final, Iterator

import pytest
import requests
from typer.testing import CliRunner, Result

from ingenious.cli import app as root_app
from ingenious.document_processing import extract as _extract_docs
from ingenious.document_processing.cli import doc_app
from ingenious.document_processing.extractor import _load

__all__ = [
    # constants
    "DEFAULT_TIMEOUT",
    # capability booleans
    "pdfminer_available",
    "unstructured_available",
    "pptx_available",
    # third-party modules / extractors
    "fitz_mod",
    "docx_mod",
    "pymupdf",
    # PDF samples (remote + synthetic)
    "pdf_path",
    "pdf_bytes",
    "sample_pdf_path",
    "sample_pdf_bytes",
    # DOCX samples
    "docx_path",
    "sample_docx_path",
    "sample_docx_bytes",
    # PPTX samples
    "pptx_path",
    "pptx_bytes",
    # cli fixtures and helpers
    "cli_runner",
    "_cli",
]

# ─────────────────── constants ───────────────────
DEFAULT_TIMEOUT: Final[int] = 10  # seconds

#: Two lightweight, publicly accessible PDFs hosted by DenseBreast-Info.
PDF_URLS: list[str] = [
    "https://www.w3.org/WAI/ER/tests/xhtml/testfiles/resources/pdf/dummy.pdf",
    "https://unec.edu.az/application/uploads/2014/12/pdf-sample.pdf",
]


# ───────────── helper functions ─────────────
def _make_one_page_pdf() -> bytes:
    """Return a minimal single-page PDF as raw bytes (pure Python)."""
    return (
        b"%PDF-1.4\n"
        b"1 0 obj<< /Type /Catalog /Pages 2 0 R >>endobj\n"
        b"2 0 obj<< /Type /Pages /Kids[3 0 R] /Count 1 >>endobj\n"
        b"3 0 obj<< /Type /Page /Parent 2 0 R /Resources<<"
        b"/Font<< /F1 4 0 R >> >> /MediaBox[0 0 595 842] /Contents 5 0 R>>endobj\n"
        b"4 0 obj<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>endobj\n"
        b"5 0 obj<< /Length 44 >>stream\n"
        b"BT /F1 12 Tf 100 700 Td (Hello PDF) Tj ET\n"
        b"endstream endobj\n"
        b"xref\n0 6\n"
        b"0000000000 65535 f \n"
        b"0000000010 00000 n \n"
        b"0000000074 00000 n \n"
        b"0000000136 00000 n \n"
        b"0000000277 00000 n \n"
        b"0000000394 00000 n \n"
        b"trailer<< /Size 6 /Root 1 0 R >>\nstartxref\n520\n%%EOF\n"
    )


def _download(url: str, cache_dir: Path) -> Path:
    """Fetch *url* into *cache_dir* with caching and fault tolerance."""
    cache_dir.mkdir(exist_ok=True)
    destination = cache_dir / Path(url).name

    if destination.exists():
        return destination

    try:
        response = requests.get(url, timeout=DEFAULT_TIMEOUT)
        response.raise_for_status()
    except requests.RequestException as exc:
        if destination.exists():
            pytest.xfail(f"Using cached PDF after download failure: {exc}")
            return destination
        pytest.skip(f"Cannot fetch {url} ({exc}); skipping network test")

    destination.write_bytes(response.content)
    return destination


def _pdf_to_text(src: Path) -> str:
    """Extract plain text from *src* using the public :pyfunc:`extract` API."""
    blocks = _extract_docs(src)
    text_lines: Iterator[str] = (block["text"] for block in blocks if block.get("text"))
    return "\n".join(text_lines).strip()


# ───────────────────── fixtures ─────────────────────
# Third-party modules -----------------------------------------------------------
@pytest.fixture(scope="session")
def fitz_mod() -> Any:
    """Import and expose the *PyMuPDF* ``fitz`` module."""
    try:
        return importlib.import_module("fitz")
    except ImportError:  # pragma: no cover – environment specific
        pytest.skip("PyMuPDF not installed")


@pytest.fixture(scope="session")
def docx_mod() -> Any:
    """Import and expose *python-docx*."""
    try:
        return importlib.import_module("docx")
    except ImportError:  # pragma: no cover
        pytest.skip("python-docx not installed")


# Capability booleans ----------------------------------------------------------
@pytest.fixture(scope="session")
def pdfminer_available() -> bool:
    """Tell whether :pypi:`pdfminer.six` is importable."""
    return importlib.util.find_spec("pdfminer") is not None


@pytest.fixture(scope="session")
def unstructured_available() -> bool:
    """Tell whether :pypi:`unstructured` is importable."""
    return importlib.util.find_spec("unstructured") is not None


@pytest.fixture(scope="session")
def pptx_available() -> bool:
    """Tell whether *python-pptx* is importable."""
    return importlib.util.find_spec("pptx") is not None


# Extractor helper -------------------------------------------------------------
@pytest.fixture(scope="session")
def pymupdf() -> Any:
    """Return a session-scoped *PyMuPDF* extractor instance."""
    return _load("pymupdf")


# PDF fixtures -----------------------------------------------------------------
@pytest.fixture(scope="session", params=PDF_URLS, ids=lambda u: Path(u).name)
def pdf_path(
    request: pytest.FixtureRequest, tmp_path_factory: pytest.TempPathFactory
) -> Path:
    """Provide each remote PDF as a cached local path."""
    cache_dir = tmp_path_factory.mktemp("pdf_cache")
    return _download(request.param, cache_dir)


@pytest.fixture()
def pdf_bytes(pdf_path: Path) -> bytes:
    """Return raw bytes of the parametrised PDF."""
    return pdf_path.read_bytes()


# DOCX fixtures ----------------------------------------------------------------
@pytest.fixture()
def docx_path(tmp_path: Path, pdf_path: Path, docx_mod: Any) -> Path:
    """Generate a DOCX mirroring the textual content of *pdf_path*."""
    destination = tmp_path / f"{pdf_path.stem}.docx"
    if destination.exists():
        return destination

    imported_text = _pdf_to_text(pdf_path) or f"Empty extract from {pdf_path.name}"
    document = docx_mod.Document()
    for line in imported_text.splitlines():
        document.add_paragraph(line)
    document.save(destination)
    return destination


# PPTX fixtures ----------------------------------------------------------------
@pytest.fixture()
def pptx_path(tmp_path: Path, pptx_available: bool) -> Path:
    """Create a minimal PPTX sample (requires *python-pptx*)."""
    if not pptx_available:
        pytest.skip("python-pptx not installed")
    from pptx import Presentation  # type: ignore

    out_path = tmp_path / "sample.pptx"
    presentation: Any = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Hello PPTX"
    presentation.save(out_path)
    return out_path


@pytest.fixture()
def pptx_bytes(pptx_path: Path) -> bytes:
    """Return raw bytes of the generated PPTX."""
    return pptx_path.read_bytes()


# Synthetic single-page PDF -----------------------------------------------------
@pytest.fixture(scope="session")
def sample_pdf_path(tmp_path_factory: pytest.TempPathFactory) -> Path:
    """Provide a simple one-page PDF stored in the session temp dir."""
    out = tmp_path_factory.getbasetemp() / "sample.pdf"
    if not out.exists():
        out.write_bytes(_make_one_page_pdf())
    return out


@pytest.fixture()
def sample_pdf_bytes(sample_pdf_path: Path) -> bytes:
    """Return raw bytes of the synthetic PDF."""
    return sample_pdf_path.read_bytes()


# DOCX generated from the synthetic PDF ----------------------------------------
@pytest.fixture(scope="session")
def sample_docx_path(
    tmp_path_factory: pytest.TempPathFactory, sample_pdf_path: Path, docx_mod: Any
) -> Path:
    """Generate a DOCX derived from the one-page synthetic PDF."""
    out = tmp_path_factory.getbasetemp() / "sample_from_pdf.docx"
    if out.exists():
        return out

    text = _pdf_to_text(sample_pdf_path) or "empty"
    document = docx_mod.Document()
    for line in text.splitlines():
        document.add_paragraph(line)
    document.save(out)
    return out


@pytest.fixture()
def sample_docx_bytes(sample_docx_path: Path) -> bytes:
    """Return raw bytes of the sample DOCX."""
    return sample_docx_path.read_bytes()


# ───────────────── CLI-level shared helpers ──────────────────
@pytest.fixture(scope="session")
def cli_runner() -> CliRunner:
    """Singleton Typer test runner reused by every test."""
    return CliRunner()


@pytest.fixture()
def _cli(cli_runner: CliRunner) -> Callable[..., Result]:
    """
    Wrapper around ``cli_runner.invoke``

    Parameters
    ----------
    source:
        Path or URL to process.
    engine:
        Name of extractor (optional).
    out_file:
        Path to output file (optional).
    entry:
        ``"sub"`` → internal Typer sub-app (``doc_app``)
        ``"root"`` → installed CLI root app (``root_app``)
    """

    def _invoke(
        source: str,
        engine: str | None = None,
        out_file: Path | None = None,
        entry: str = "sub",
    ) -> Result:
        args: list[str] = []
        if entry == "root":
            args.append("document-processing")
            target = root_app
        else:
            target = doc_app

        args.append(source)
        if engine is not None:
            args.extend(["--engine", engine])
        if out_file is not None:
            args.extend(["--out", str(out_file)])

        return cli_runner.invoke(target, args)

    return _invoke
